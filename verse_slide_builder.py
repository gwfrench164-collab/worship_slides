from pathlib import Path
import re
import math

from pptx_utils import (
    load_template,
    find_template_slide_index,
    add_scripture_slide_from_template,
    TOKEN_VERSE_REF,
    TOKEN_VERSE_TXT,
)

from pptx.util import Pt

def _emu_to_points(emu: int) -> float:
    # 1 point = 12700 EMU
    return emu / 12700.0

def find_shape_with_token(slide, token: str):
    for shape in slide.shapes:
        if shape.has_text_frame and token in shape.text:
            return shape
    return None

def estimate_text_capacity_from_shape(shape) -> int:
    """
    Estimate how many characters fit in this textbox using its width/height
    and the template font size. This adapts automatically when the user
    changes the font size in the template.
    """
    width_pts = _emu_to_points(shape.width)
    height_pts = _emu_to_points(shape.height)

    tf = shape.text_frame
    p0 = tf.paragraphs[0]
    font = p0.runs[0].font if p0.runs else p0.font

    # fallback if template doesn’t have size set
    font_size_pts = font.size.pt if font.size else 44.0

    # These constants are “reasonable worship-slide heuristics”
    avg_char_w = font_size_pts * 0.43     # character width estimate
    line_h = font_size_pts * 1.15         # line height estimate

    chars_per_line = max(10, int(width_pts / max(avg_char_w, 1.0)))
    lines_in_box = max(1, int(height_pts / max(line_h, 1.0)))

    # Safety margin so we don’t overflow
    max_chars = int(chars_per_line * lines_in_box * 0.90)
    return max(80, max_chars)  # never go absurdly low
# -----------------------
# Bracket-span protection
# -----------------------
_NBSP = "\u00A0"
_BRACKET_SPAN_RE = re.compile(r"\[(.+?)\]")

def protect_bracket_spans_for_wrapping(text: str) -> str:
    """
    Prevent wrapping from splitting inside [bracketed spans] by converting
    spaces inside the brackets into NBSP (non-breaking spaces).
    """
    def repl(m: re.Match) -> str:
        inner = " ".join(m.group(1).split())
        inner = inner.replace(" ", _NBSP)
        return f"[{inner}]"
    return _BRACKET_SPAN_RE.sub(repl, text)

def restore_spaces(text: str) -> str:
    return text.replace(_NBSP, " ")

# -----------------------
# Smart wrap (punctuation friendly)
# -----------------------
_SMALL_WORDS = {
    "and","or","but","the","a","an","of","to","in","on","at","for","by","with",
    "that","this","is","be","as","if","so","yet","nor","from"
}
_PUNCT_END = re.compile(r"[,:;.!?]$")

def wrap_verse_smart(text: str, width: int) -> list[str]:
    text = " ".join(text.split())
    if not text:
        return []

    words = text.split(" ")
    lines = []
    i = 0

    while i < len(words):
        j = i
        line = words[j]
        j += 1
        while j < len(words) and len(line) + 1 + len(words[j]) <= width:
            line += " " + words[j]
            j += 1

        # Avoid ending a line with a tiny word
        parts = line.split(" ")
        if len(parts) >= 2 and parts[-1].lower() in _SMALL_WORDS and j < len(words):
            parts.pop()
            line = " ".join(parts)
            j -= 1

        lines.append(line)
        i = j

    return lines

# -----------------------
# FIT MODE: measure template box and font
# -----------------------
EMU_PER_INCH = 914400
PT_PER_INCH = 72

def emu_to_pt(emu: int) -> float:
    return (emu / EMU_PER_INCH) * PT_PER_INCH

def _best_font_size_pts_from_shape(shape) -> float:
    """
    Try hard to get font size from the token shape. Keynote exports often put
    the "real" formatting on the first run.
    """
    try:
        tf = shape.text_frame
        if not tf.paragraphs:
            return 60.0
        p0 = tf.paragraphs[0]

        # Prefer first run font size
        if p0.runs and p0.runs[0].font.size:
            return float(p0.runs[0].font.size.pt)

        # Fallback to paragraph font
        if p0.font and p0.font.size:
            return float(p0.font.size.pt)

    except Exception:
        pass

    return 60.0  # safe fallback for your case

def _line_spacing_factor_from_shape(shape, font_size_pts: float) -> float:
    """
    Convert line spacing info to a multiplier. Keynote templates often use defaults.
    """
    try:
        p0 = shape.text_frame.paragraphs[0]
        ls = p0.line_spacing
        if ls is None:
            return 1.10  # default-ish
        # If it's a float (like 1.0, 1.15, 1.2)
        if isinstance(ls, (float, int)):
            # if it's "points" (big number), convert to multiplier
            if ls > 3:
                return float(ls) / max(font_size_pts, 1.0)
            return float(ls)
        # If it's a pptx Length object with .pt
        if hasattr(ls, "pt"):
            return float(ls.pt) / max(font_size_pts, 1.0)
    except Exception:
        pass

    return 1.10

def _find_token_shape(slide, token: str):
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        try:
            if token in shape.text:
                return shape
        except Exception:
            continue
    return None

def compute_fit_params(prs, tpl_idx: int, token: str, preset: str = "normal") -> tuple[int, int]:
    """
    Returns (chars_per_line, max_lines_per_slide) estimated from the template box.
    preset: tight / normal / loose
    """
    slide = prs.slides[tpl_idx]
    shape = _find_token_shape(slide, token)
    if shape is None:
        # fallback
        return (34, 4)

    width_pts = emu_to_pt(int(shape.width))
    height_pts = emu_to_pt(int(shape.height))

    font_size_pts = _best_font_size_pts_from_shape(shape)
    line_factor = _line_spacing_factor_from_shape(shape, font_size_pts)

    # 60pt worship fonts usually render closer to ~0.40–0.46 of font size per character.
    # 0.43 is a good default; it will use more of the textbox width and wrap later.
    avg_char_w = font_size_pts * 0.43
    raw_cpl = max(10, int(width_pts / max(avg_char_w, 1.0)))

    # Max lines is textbox height divided by line height
    line_height_pts = font_size_pts * line_factor
    raw_lines = max(1, int(height_pts / max(line_height_pts, 1.0)))

    # Presets (A preference: readability / natural phrasing)
    preset = (preset or "normal").lower().strip()
    if preset == "tight":
        cpl = max(10, int(raw_cpl * 0.92))
        max_lines = max(1, raw_lines - 1)
    elif preset == "loose":
        cpl = max(10, int(raw_cpl * 1.05))
        max_lines = raw_lines
    else:  # normal
        cpl = raw_cpl
        max_lines = raw_lines

    # Guardrails for 60pt: don’t let it get too ambitious
    # Allow wider lines; your textbox is very wide.
    cpl = min(cpl, 60)
    max_lines = min(max_lines, 6)

    return cpl, max_lines

def group_lines_into_slides_phrasey(lines: list[str], max_lines: int) -> list[list[str]]:
    """
    Group wrapped lines into slide chunks, preferring to break after punctuation.
    """
    if not lines:
        return []

    slides = []
    i = 0
    while i < len(lines):
        chunk = lines[i:i+max_lines]

        # If this chunk ends mid-sentence and we have room to shift, try to break earlier
        if i + max_lines < len(lines):
            if chunk and not _PUNCT_END.search(chunk[-1].strip()):
                # Look backward in the chunk for a better break point
                best = None
                for k in range(len(chunk)-1, -1, -1):
                    if _PUNCT_END.search(chunk[k].strip()):
                        best = k + 1
                        break
                # If we found punctuation break and it's not too tiny, use it
                if best is not None and best >= max(2, max_lines - 2):
                    chunk = lines[i:i+best]
                    i += best
                    slides.append(chunk)
                    continue

        slides.append(chunk)
        i += max_lines

    return slides

# -----------------------
# Public entry
# -----------------------
def build_verse_deck(
    template_path: Path,
    refs_and_texts: list[tuple[str, str]],
    output_path: Path,
    fit_preset: str = "normal",   # tight / normal / loose
):
    prs = load_template(template_path)

    # Find the scripture template slide by tokens (user can move it anywhere)
    tpl_idx = find_template_slide_index(prs, [TOKEN_VERSE_REF, TOKEN_VERSE_TXT])

    tpl_slide = prs.slides[tpl_idx]
    verse_box = find_shape_with_token(tpl_slide, TOKEN_VERSE_TXT)
    if verse_box is None:
        raise RuntimeError("Could not find verse text box containing {{VERSE TXT}} on template slide.")

    max_chars = estimate_text_capacity_from_shape(verse_box)

    # Compute wrapping targets from actual template textbox
    chars_per_line, max_lines = compute_fit_params(prs, tpl_idx, TOKEN_VERSE_TXT, preset=fit_preset)

    for ref, verse_text in refs_and_texts:
        verse_text = " ".join((verse_text or "").split())
        if not verse_text:
            add_scripture_slide_from_template(prs, tpl_idx, ref, "(text unavailable)")
            continue

        # Protect bracketed spans so we never split inside them
        protected = protect_bracket_spans_for_wrapping(verse_text)

        # Wrap based on computed chars-per-line
        wrapped_lines = wrap_verse_smart(protected, width=chars_per_line)

        # Restore normal spaces for rendering + italics conversion
        wrapped_lines = [restore_spaces(l) for l in wrapped_lines]

        # Group lines into slides, preferring punctuation breaks
        slide_chunks = group_lines_into_slides_phrasey(wrapped_lines, max_lines=max_lines)

        for chunk_lines in slide_chunks:
            # We pass a SINGLE string with line breaks to preserve “line intent”
            # while still letting the template style control appearance.
            slide_text = "\n".join(chunk_lines)
            add_scripture_slide_from_template(prs, tpl_idx, ref, slide_text)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)