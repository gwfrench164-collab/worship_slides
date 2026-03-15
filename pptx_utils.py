from copy import deepcopy
from pathlib import Path
from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
import re
import posixpath
import shutil
import tempfile
import zipfile
import xml.etree.ElementTree as ET


def load_template(path):
    return Presentation(path)


# -------------------------
# Slide duplication (preserves Keynote/PowerPoint styling)
# -------------------------

def _copy_relationships(src_slide, dst_slide):
    src_part = src_slide.part
    dst_part = dst_slide.part

    for _rId, rel_obj in src_part.rels.items():
        if rel_obj.is_external:
            continue
        if rel_obj.reltype == RT.SLIDE_LAYOUT:
            continue
        dst_part.rels._add_relationship(rel_obj.reltype, rel_obj.target_part)


def duplicate_slide(prs, slide_index: int):
    src = prs.slides[slide_index]
    blank_layout = prs.slide_layouts[0]
    dst = prs.slides.add_slide(blank_layout)

    # remove default shapes
    for shape in list(dst.shapes):
        el = shape._element
        el.getparent().remove(el)

    # copy background
    if src._element.bg is not None and len(src._element.bg) > 0:
        dst._element.get_or_add_bg()
        dst._element.bg.clear()
        dst._element.bg.append(deepcopy(src._element.bg[0]))

    # copy shapes
    for shape in src.shapes:
        new_el = deepcopy(shape._element)
        dst.shapes._spTree.insert_element_before(new_el, "p:extLst")

    _copy_relationships(src, dst)
    return dst


# ---- ZIP/XML-based cross-presentation merge helpers ----

P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
CT_NS = "http://schemas.openxmlformats.org/package/2006/content-types"

ET.register_namespace("", P_NS)
ET.register_namespace("a", "http://schemas.openxmlformats.org/drawingml/2006/main")
ET.register_namespace("r", R_NS)


def _ns(tag, namespace):
    return f"{{{namespace}}}{tag}"


def _slide_rel_path(slide_num: int) -> str:
    return f"ppt/slides/_rels/slide{slide_num}.xml.rels"


def _slide_xml_path(slide_num: int) -> str:
    return f"ppt/slides/slide{slide_num}.xml"


def _next_part_number(names: list[str], prefix: str, suffix: str) -> int:
    max_num = 0
    for name in names:
        if name.startswith(prefix) and name.endswith(suffix):
            middle = name[len(prefix):-len(suffix)]
            if middle.isdigit():
                max_num = max(max_num, int(middle))
    return max_num + 1


def _copy_tree(src: Path, dst: Path):
    if dst.exists():
        shutil.rmtree(dst)
    shutil.copytree(src, dst)


def _unzip_to_dir(pptx_path: Path, out_dir: Path):
    with zipfile.ZipFile(pptx_path, "r") as zf:
        zf.extractall(out_dir)


def _zip_dir_to_pptx(src_dir: Path, pptx_path: Path):
    with zipfile.ZipFile(pptx_path, "w", zipfile.ZIP_DEFLATED) as zf:
        for file_path in sorted(src_dir.rglob("*")):
            if file_path.is_file():
                arcname = file_path.relative_to(src_dir).as_posix()
                zf.write(file_path, arcname)


def _read_xml(path: Path) -> ET.ElementTree:
    return ET.parse(path)


def _write_xml(tree: ET.ElementTree, path: Path):
    tree.write(path, encoding="UTF-8", xml_declaration=True)


def _collect_slide_numbers(root_dir: Path) -> list[int]:
    slides_dir = root_dir / "ppt" / "slides"
    nums = []
    for path in slides_dir.glob("slide*.xml"):
        stem = path.stem
        num = stem.replace("slide", "")
        if num.isdigit():
            nums.append(int(num))
    return sorted(nums)


def _collect_media_names(root_dir: Path) -> list[str]:
    media_dir = root_dir / "ppt" / "media"
    if not media_dir.exists():
        return []
    return sorted(p.name for p in media_dir.iterdir() if p.is_file())


def _ensure_content_type_override(content_types_root: ET.Element, part_name: str, content_type: str):
    for child in content_types_root.findall(_ns("Override", CT_NS)):
        if child.attrib.get("PartName") == part_name:
            return
    ET.SubElement(
        content_types_root,
        _ns("Override", CT_NS),
        {"PartName": part_name, "ContentType": content_type},
    )


def _add_slide_to_presentation_xml(presentation_root: ET.Element, slide_rid: str):
    sld_id_lst = presentation_root.find(_ns("sldIdLst", P_NS))
    if sld_id_lst is None:
        sld_id_lst = ET.SubElement(presentation_root, _ns("sldIdLst", P_NS))

    existing_ids = []
    for child in sld_id_lst.findall(_ns("sldId", P_NS)):
        try:
            existing_ids.append(int(child.attrib.get("id", "0")))
        except ValueError:
            pass

    next_id = max(existing_ids, default=255) + 1
    ET.SubElement(
        sld_id_lst,
        _ns("sldId", P_NS),
        {"id": str(next_id), _ns("id", R_NS): slide_rid},
    )


def _add_relationship(rel_root: ET.Element, rel_type: str, target: str) -> str:
    existing = []
    for rel in rel_root.findall(_ns("Relationship", REL_NS)):
        rid = rel.attrib.get("Id", "")
        if rid.startswith("rId") and rid[3:].isdigit():
            existing.append(int(rid[3:]))

    next_rid = f"rId{max(existing, default=0) + 1}"
    ET.SubElement(
        rel_root,
        _ns("Relationship", REL_NS),
        {"Id": next_rid, "Type": rel_type, "Target": target},
    )
    return next_rid


def _remap_slide_relationship_targets(src_root: Path, dst_root: Path, src_rel_path: str, dst_rel_path: str):
    src_rel_full = src_root / src_rel_path
    if not src_rel_full.exists():
        return

    dst_rel_full = dst_root / dst_rel_path
    dst_rel_full.parent.mkdir(parents=True, exist_ok=True)

    rel_tree = _read_xml(src_rel_full)
    rel_root = rel_tree.getroot()

    dst_media_dir = dst_root / "ppt" / "media"
    dst_media_dir.mkdir(parents=True, exist_ok=True)

    existing_media = _collect_media_names(dst_root)

    for rel in rel_root.findall(_ns("Relationship", REL_NS)):
        target = rel.attrib.get("Target", "")
        if not target.startswith("../media/"):
            continue

        old_name = posixpath.basename(target)
        prefix, dot, ext = old_name.partition(".")
        if not dot:
            continue

        next_num = _next_part_number(existing_media, prefix.rstrip("0123456789") or "image", f".{ext}")
        base_prefix = re.sub(r"\d+$", "", prefix) or "image"
        new_name = f"{base_prefix}{next_num}.{ext}"

        shutil.copy2(src_root / "ppt" / "media" / old_name, dst_media_dir / new_name)
        existing_media.append(new_name)
        rel.set("Target", f"../media/{new_name}")

    _write_xml(rel_tree, dst_rel_full)


def merge_presentations(song_deck_path, verse_deck_path, output_path):
    song_deck_path = Path(song_deck_path)
    verse_deck_path = Path(verse_deck_path)
    output_path = Path(output_path)

    song_prs = Presentation(song_deck_path)
    verse_prs = Presentation(verse_deck_path)

    if (
        song_prs.slide_width != verse_prs.slide_width
        or song_prs.slide_height != verse_prs.slide_height
    ):
        raise RuntimeError(
            "The song deck and verse deck do not use the same slide size. Build both decks from templates with matching dimensions before merging."
        )

    with tempfile.TemporaryDirectory() as tmpdir:
        tmpdir = Path(tmpdir)
        song_dir = tmpdir / "song"
        verse_dir = tmpdir / "verse"
        merged_dir = tmpdir / "merged"

        _unzip_to_dir(song_deck_path, song_dir)
        _unzip_to_dir(verse_deck_path, verse_dir)
        _copy_tree(song_dir, merged_dir)

        merged_slide_nums = _collect_slide_numbers(merged_dir)
        next_slide_num = max(merged_slide_nums, default=0) + 1

        presentation_tree = _read_xml(merged_dir / "ppt" / "presentation.xml")
        presentation_root = presentation_tree.getroot()

        pres_rels_tree = _read_xml(merged_dir / "ppt" / "_rels" / "presentation.xml.rels")
        pres_rels_root = pres_rels_tree.getroot()

        content_types_tree = _read_xml(merged_dir / "[Content_Types].xml")
        content_types_root = content_types_tree.getroot()

        for verse_slide_num in _collect_slide_numbers(verse_dir):
            new_slide_num = next_slide_num
            next_slide_num += 1

            src_slide_xml = verse_dir / _slide_xml_path(verse_slide_num)
            dst_slide_xml = merged_dir / _slide_xml_path(new_slide_num)
            dst_slide_xml.parent.mkdir(parents=True, exist_ok=True)
            shutil.copy2(src_slide_xml, dst_slide_xml)

            _remap_slide_relationship_targets(
                verse_dir,
                merged_dir,
                _slide_rel_path(verse_slide_num),
                _slide_rel_path(new_slide_num),
            )

            new_rid = _add_relationship(
                pres_rels_root,
                "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide",
                f"slides/slide{new_slide_num}.xml",
            )
            _add_slide_to_presentation_xml(presentation_root, new_rid)
            _ensure_content_type_override(
                content_types_root,
                f"/ppt/slides/slide{new_slide_num}.xml",
                "application/vnd.openxmlformats-officedocument.presentationml.slide+xml",
            )

        _write_xml(presentation_tree, merged_dir / "ppt" / "presentation.xml")
        _write_xml(pres_rels_tree, merged_dir / "ppt" / "_rels" / "presentation.xml.rels")
        _write_xml(content_types_tree, merged_dir / "[Content_Types].xml")

        output_path.parent.mkdir(parents=True, exist_ok=True)
        _zip_dir_to_pptx(merged_dir, output_path)

def remove_slide(prs, index: int):
    slide_id_list = prs.slides._sldIdLst  # pylint: disable=protected-access
    slides = list(slide_id_list)
    slide_id_list.remove(slides[index])


# -------------------------
# Token finding + replacement
# -------------------------

def _slide_text_contains(slide, token: str) -> bool:
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        try:
            if token in shape.text:
                return True
        except Exception:
            pass
    return False


def find_template_slide_index(prs, required_tokens: list[str]) -> int:
    for i, slide in enumerate(prs.slides):
        if all(_slide_text_contains(slide, tok) for tok in required_tokens):
            return i
    raise RuntimeError(f"Template slide not found containing tokens: {required_tokens}")


def _get_best_font_source(paragraph):
    if paragraph.runs:
        return paragraph.runs[0].font
    return paragraph.font


def _copy_font_style(dst_font, src_font):
    dst_font.name = src_font.name
    dst_font.size = src_font.size
    dst_font.bold = src_font.bold
    dst_font.italic = src_font.italic

    if src_font.color is not None:
        try:
            if src_font.color.type == 1 and src_font.color.rgb is not None:
                dst_font.color.rgb = src_font.color.rgb
            elif src_font.color.type == 2 and src_font.color.theme_color is not None:
                dst_font.color.theme_color = src_font.color.theme_color
        except Exception:
            pass


def _force_alignment_like_template(p, template_alignment):
    if template_alignment is None or template_alignment == PP_ALIGN.LEFT:
        p.alignment = PP_ALIGN.LEFT
    else:
        p.alignment = template_alignment


def _replace_token_text(slide, token: str, new_text: str) -> bool:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        tf = shape.text_frame
        if token not in tf.text:
            continue

        p0 = tf.paragraphs[0]
        alignment = p0.alignment
        level = p0.level
        space_before = p0.space_before
        space_after = p0.space_after
        line_spacing = p0.line_spacing
        src_font = _get_best_font_source(p0)

        lines = (new_text or "").split("\n")
        tf.clear()

        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line

            _force_alignment_like_template(p, alignment)
            p.level = level
            p.space_before = space_before
            p.space_after = space_after
            p.line_spacing = line_spacing

            _copy_font_style(p.font, src_font)

        return True

    return False


_BRACKET_ITALIC_RE = re.compile(r"\[(.+?)\]")


def _replace_token_text_with_bracket_italics(slide, token: str, new_text: str) -> bool:
    """
    Replace token text while italicizing bracketed spans like [this].
    Brackets are removed from the final slide text.

    This implementation is *streaming* across the whole text so bracket spans
    remain italicized even if wrapping/splitting introduces newlines inside a
    bracketed span.
    """
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        tf = shape.text_frame
        if token not in tf.text:
            continue

        p0 = tf.paragraphs[0]
        alignment = p0.alignment
        level = p0.level
        space_before = p0.space_before
        space_after = p0.space_after
        line_spacing = p0.line_spacing
        src_font = _get_best_font_source(p0)

        text = new_text or ""

        # Clear and rebuild paragraphs/runs, copying style from the template run/font.
        tf.clear()

        def setup_paragraph(p):
            _force_alignment_like_template(p, alignment)
            p.level = level
            p.space_before = space_before
            p.space_after = space_after
            p.line_spacing = line_spacing

        # Streaming parse: remove '[' and ']' and toggle italic state.
        italic = False
        cur_para = tf.paragraphs[0]
        setup_paragraph(cur_para)
        cur_para.text = ""  # ensure empty

        def add_run(p, s, ital):
            if not s:
                return
            r = p.add_run()
            r.text = s
            _copy_font_style(r.font, src_font)
            r.font.italic = bool(ital)

        buf = []
        buf_ital = False

        def flush_buf(p):
            nonlocal buf, buf_ital
            if buf:
                add_run(p, "".join(buf), buf_ital)
                buf = []

        # Initialize buffer state
        buf_ital = italic

        for ch in text:
            if ch == "[":
                # Flush current buffer before toggling
                flush_buf(cur_para)
                italic = True
                buf_ital = italic
                continue
            if ch == "]":
                flush_buf(cur_para)
                italic = False
                buf_ital = italic
                continue
            if ch == "\n":
                flush_buf(cur_para)
                # New paragraph
                cur_para = tf.add_paragraph()
                setup_paragraph(cur_para)
                cur_para.text = ""
                buf_ital = italic
                continue

            # Normal character
            # If italic state changes mid-buffer (shouldn't happen), flush
            if italic != buf_ital:
                flush_buf(cur_para)
                buf_ital = italic
            buf.append(ch)

        flush_buf(cur_para)
        return True

    return False


# -------------------------
# Public token helpers
# -------------------------

TOKEN_TITLE = "{{TITLE}}"
TOKEN_AUTHOR = "{{AUTHOR}}"
TOKEN_CCLI = "{{CCLI}}"
TOKEN_LYRICS = "{{LYRICS}}"
TOKEN_VERSE_REF = "{{VERSE REF}}"
TOKEN_VERSE_TXT = "{{VERSE TXT}}"


def add_title_slide_from_template(
    prs,
    template_slide_index: int,
    title_text: str,
    author_text: str = "",
    ccli_text: str = "",
):
    slide = duplicate_slide(prs, template_slide_index)
    if not _replace_token_text(slide, TOKEN_TITLE, title_text):
        raise RuntimeError("Template title slide missing {{TITLE}} token.")

    # These are optional. If the template contains the token, replace it.
    # If the value is blank, leave the textbox blank.
    _replace_token_text(slide, TOKEN_AUTHOR, author_text or "")
    _replace_token_text(slide, TOKEN_CCLI, ccli_text or "")
    return slide

from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.util import Inches


def add_lyrics_slide_from_template(
    prs,
    template_slide_index: int,
    lines: list[str],
    *,
    lyric_starts: list[bool] | None = None,
    lyric_gap_pt: float = 0.0,
    hanging_indent_pt: float = 0.0,
):
    """
    Adds a lyrics slide from the template.

    `lines` are DISPLAY lines (already wrapped). We keep each display line as its own paragraph.
    `lyric_starts[i] == True` means lines[i] is the first display line of a NEW lyric line,
    so we add `space_before` (instead of inserting a blank paragraph that wastes a whole line).

    This keeps visual separation between lyric lines while allowing better slide packing.
    """
    slide = duplicate_slide(prs, template_slide_index)
    lyric_text = "\n".join(lines)

    # Replace token
    if not _replace_token_text(slide, TOKEN_LYRICS, lyric_text):
        raise RuntimeError("Template lyrics slide missing {{LYRICS}} token.")

    # Find the lyrics shape (prefer exact name, then best-effort fallback)
    lyrics_shape = None
    for sh in slide.shapes:
        try:
            if getattr(sh, "name", None) == TOKEN_LYRICS:
                lyrics_shape = sh
                break
        except Exception:
            pass

    if lyrics_shape is None:
        # Fallback: first non-empty text frame
        for sh in slide.shapes:
            if getattr(sh, "has_text_frame", False):
                try:
                    if sh.text_frame and sh.text_frame.text.strip():
                        lyrics_shape = sh
                        break
                except Exception:
                    pass

    if lyrics_shape is None or not getattr(lyrics_shape, "has_text_frame", False):
        return slide

    tf = lyrics_shape.text_frame

    # Ensure PowerPoint does NOT re-wrap our manually wrapped lines
    try:
        tf.word_wrap = False
    except Exception:
        pass

    # Prevent auto-sizing from changing our layout
    try:
        tf.auto_size = MSO_AUTO_SIZE.NONE
    except Exception:
        try:
            tf.auto_size = None
        except Exception:
            pass

    # Apply paragraph spacing + optional hanging indent
    flags = lyric_starts if (lyric_starts and len(lyric_starts) == len(tf.paragraphs)) else None
    for i, p in enumerate(tf.paragraphs):
        # Paragraph spacing for lyric separation (no blank lines)
        if flags and i > 0 and flags[i] and lyric_gap_pt and lyric_gap_pt > 0:
            p.space_before = Pt(float(lyric_gap_pt))

        # Optional hanging indent (kept for compatibility)
        if hanging_indent_pt and hanging_indent_pt > 0:
            p.left_indent = Pt(hanging_indent_pt)
            p.first_line_indent = Pt(-hanging_indent_pt)

    return slide
def add_scripture_slide_from_template(prs, template_slide_index: int, verse_ref: str, verse_text: str):
    slide = duplicate_slide(prs, template_slide_index)

    # Capture the verse textbox shape BEFORE replacement so we can apply guardrails reliably.
    verse_shape = None
    ref_shape = None
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        tf = shape.text_frame
        try:
            txt = tf.text or ""
        except Exception:
            txt = ""
        if TOKEN_VERSE_TXT in txt:
            verse_shape = shape
        if TOKEN_VERSE_REF in txt:
            ref_shape = shape

    ok1 = _replace_token_text(slide, TOKEN_VERSE_REF, verse_ref)
    ok2 = _replace_token_text_with_bracket_italics(slide, TOKEN_VERSE_TXT, verse_text)

    if not ok1:
        raise RuntimeError("Scripture template slide missing {{VERSE REF}} token.")
    if not ok2:
        raise RuntimeError("Scripture template slide missing {{VERSE TXT}} token.")

    # --- Guardrails: prevent PowerPoint from 'helping' in surprising ways ---
    # We do manual line breaks for verses; avoid autoshrink and unintended reflow.
    if verse_shape is not None and getattr(verse_shape, "has_text_frame", False):
        tf = verse_shape.text_frame
        try:
            tf.word_wrap = False
        except Exception:
            pass
        try:
            tf.auto_size = MSO_AUTO_SIZE.NONE
        except Exception:
            try:
                tf.auto_size = None
            except Exception:
                pass

    # Keep the reference stable too (no autoshrink surprises).
    if ref_shape is not None and getattr(ref_shape, "has_text_frame", False):
        tf = ref_shape.text_frame
        try:
            tf.word_wrap = False
        except Exception:
            pass
        try:
            tf.auto_size = MSO_AUTO_SIZE.NONE
        except Exception:
            try:
                tf.auto_size = None
            except Exception:
                pass

    return slide


def add_debug_guides(slide, target_shape, *, usable_rect_emu=None, caption: str = ""):
    """Draw debug overlays.

    usable_rect_emu: (left, top, width, height) in EMU of the usable text area
    (textbox minus margins). If None, uses target_shape's box.
    """
    try:
        if usable_rect_emu is None:
            l, t, w, h = int(target_shape.left), int(target_shape.top), int(target_shape.width), int(target_shape.height)
        else:
            l, t, w, h = map(int, usable_rect_emu)
        # Outline rectangle, no fill
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
        rect.fill.background()  # transparent
        rect.line.width = Pt(1)
        # Keep default line color (don't force theme colors)
    except Exception:
        pass
    if caption:
        try:
            # small textbox in top-left corner
            cap = slide.shapes.add_textbox(int(target_shape.left), int(target_shape.top) - int(0.35 * 914400), int(target_shape.width), int(0.35 * 914400))
            tf = cap.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = caption
            p.font.size = Pt(10)
        except Exception:
            pass
