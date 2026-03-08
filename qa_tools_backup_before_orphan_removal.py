from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple
import re

from pptx import Presentation

EMU_PER_PT = 12700

_CONJ_START_RE = re.compile(r"^(and|but|so|for|yet|or|nor)\b", re.IGNORECASE)

def _iter_text_shapes(slide):
    for sh in slide.shapes:
        if getattr(sh, "has_text_frame", False):
            try:
                txt = sh.text_frame.text
            except Exception:
                txt = ""
            if txt is not None:
                yield sh, txt

def _norm(s: str) -> str:
    return " ".join((s or "").replace("\r", "\n").split())

def _split_lines(s: str) -> List[str]:
    return [ln.strip() for ln in (s or "").replace("\r", "\n").split("\n") if ln.strip()]

def _ends_sentence(s: str) -> bool:
    s = (s or "").strip()
    return bool(s) and s[-1] in ".!?:;\"”’)"  # commas don't count as "done"

def _is_token_slide(full_text: str) -> bool:
    return ("{{" in full_text) and ("}}" in full_text)

def _min_font_pt(slide) -> Optional[float]:
    m: Optional[float] = None
    for sh, _ in _iter_text_shapes(slide):
        try:
            for p in sh.text_frame.paragraphs:
                # first run is usually representative
                if p.runs and p.runs[0].font.size:
                    pt = float(p.runs[0].font.size.pt)
                elif p.font and p.font.size:
                    pt = float(p.font.size.pt)
                else:
                    continue
                m = pt if m is None else min(m, pt)
        except Exception:
            continue
    return m

@dataclass
class SlideText:
    slide: int
    raw_texts: List[str]
    full: str
    lines: List[str]
    chars: int
    is_token: bool
    is_title: bool
    song_title: Optional[str]
    min_font_pt: Optional[float]

def _classify_song_title(slide: int, texts: List[str]) -> Tuple[bool, Optional[str]]:
    """
    Heuristic: Title slide usually has exactly one non-empty text shape (title),
    and no multi-line lyric block.
    """
    cleaned = [t.strip() for t in texts if t and t.strip()]
    if not cleaned:
        return False, None
    # If there's one block and it's relatively short, treat as title.
    if len(cleaned) == 1:
        t = _norm(cleaned[0])
        # Titles are usually short; allow punctuation/numbers.
        if 0 < len(t) <= 80 and len(t.split()) <= 10:
            return True, t
    # If there are two blocks and one is very short (title) and the other is empty-ish, also title.
    if len(cleaned) == 2:
        a, b = map(_norm, cleaned)
        short = a if len(a) <= len(b) else b
        long = b if short == a else a
        if len(short) <= 80 and len(short.split()) <= 10 and len(long) <= 5:
            return True, short
    return False, None

def _extract_slide_text(prs: Presentation, idx: int, current_song: Optional[str]) -> SlideText:
    slide = prs.slides[idx-1]
    texts: List[str] = []
    for _, txt in _iter_text_shapes(slide):
        t = (txt or "").strip()
        if t:
            texts.append(t)

    full = "\n".join(texts)
    is_token = _is_token_slide(full)
    lines = _split_lines(full)
    chars = len(_norm(full))
    is_title, title = _classify_song_title(idx, texts)

    song_title = current_song
    if is_title and title:
        song_title = title

    return SlideText(
        slide=idx,
        raw_texts=texts,
        full=full,
        lines=lines,
        chars=chars,
        is_token=is_token,
        is_title=is_title,
        song_title=song_title,
        min_font_pt=_min_font_pt(slide),
    )

def analyze_pptx(pptx_path: Path) -> dict:
    """
    QA heuristics for quickly spotting bad slides.

    This returns:
      - flags: SPARSE/CROWDED/TINY_TEXT and (for song decks) TAIL/ORPHAN_START
      - slides: per-slide metrics
      - issues: rich issue bundle (slide text + prev/next + song title)
    """
    prs = Presentation(str(pptx_path))

    sparse: List[int] = []
    crowded: List[int] = []
    tiny: List[int] = []
    tail: List[int] = []
    orphan_start: List[int] = []

    slide_stats: List[Dict[str, Any]] = []
    issues: List[Dict[str, Any]] = []

    # Track current song based on detected title slides
    current_song: Optional[str] = None

    # Pre-extract all slide texts so we can look at prev/next easily
    slides: List[SlideText] = []
    for i in range(1, len(prs.slides) + 1):
        st = _extract_slide_text(prs, i, current_song)
        if st.is_title and st.song_title:
            current_song = st.song_title
            # refresh with the updated current_song (title itself)
            st.song_title = current_song
        slides.append(st)

    # Determine if this looks like a song deck: multiple title slides detected
    title_count = sum(1 for s in slides if s.is_title and not s.is_token)
    looks_like_song_deck = title_count >= 2

    for st in slides:
        # Ignore token-only template slides
        if st.is_token:
            slide_stats.append(
                {"slide": st.slide, "ignored": True, "chars": 0, "lines": 0, "min_font_pt": st.min_font_pt}
            )
            continue

        slide_stats.append(
            {
                "slide": st.slide,
                "chars": st.chars,
                "lines": len(st.lines),
                "min_font_pt": st.min_font_pt,
                "song_title": st.song_title,
                "is_title": st.is_title,
            }
        )

        # Skip title slides from sparse/tail/orphan checks (titles are expected to be sparse).
        if st.is_title and looks_like_song_deck:
            continue

        # Conservative crowded
        if st.chars > 380 or len(st.lines) > 9:
            crowded.append(st.slide)

        # Sparse heuristic
        if st.chars > 0 and (st.chars < 45 or len(st.lines) < 2):
            sparse.append(st.slide)

        # Tiny text heuristic (autosize)
        if st.min_font_pt is not None and st.min_font_pt < 26:
            tiny.append(st.slide)

    # Song-specific flags: TAIL + ORPHAN_START
    if looks_like_song_deck:
        for i, st in enumerate(slides):
            if st.is_token or st.is_title:
                continue
            # Tail: very short leftover fragment that likely should have been merged
            if st.chars > 0:

                prev = slides[i - 1] if i > 0 else None

                is_short = len(st.lines) <= 2
                very_short_text = st.chars < 40

                # only consider it a tail if the previous slide looks like a continuation
                continuation = False
                if prev and not prev.is_title:
                    prev_lines = _split_lines(prev.full)
                    if prev_lines:
                        last_prev = prev_lines[-1]
                        continuation = not _ends_sentence(last_prev)

                if is_short and very_short_text and continuation:
                    tail.append(st.slide)

            # Orphan start: begins with conjunction and prior slide didn't "complete" a sentence.
            if st.lines:
                first = st.lines[0]
                if _CONJ_START_RE.match(first):
                    prev = slides[i-1] if i > 0 else None
                    prev_text = prev.full if prev else ""
                    prev_last_line = _split_lines(prev_text)[-1] if prev and _split_lines(prev_text) else ""
                    if prev and (not prev.is_title) and (not _ends_sentence(prev_last_line)):
                        orphan_start.append(st.slide)

    # Build rich issue bundle (only for slides that were flagged)
    flagged_set = set(sparse) | set(crowded) | set(tiny) | set(tail) | set(orphan_start)

    for st in slides:
        if st.slide not in flagged_set:
            continue
        prev = slides[st.slide-2] if st.slide >= 2 else None
        nxt = slides[st.slide] if st.slide < len(slides) else None

        issue_types: List[str] = []
        if st.slide in sparse: issue_types.append("SPARSE")
        if st.slide in crowded: issue_types.append("CROWDED")
        if st.slide in tiny: issue_types.append("TINY_TEXT")
        if st.slide in tail: issue_types.append("TAIL")
        if st.slide in orphan_start: issue_types.append("ORPHAN_START")

        issues.append(
            {
                "slide": st.slide,
                "song_title": st.song_title,
                "is_title": st.is_title,
                "types": issue_types,
                "text": st.full,
                "prev_text": prev.full if prev else "",
                "next_text": nxt.full if nxt else "",
            }
        )

    return {
        "pptx": str(pptx_path),
        "slide_count": len(prs.slides),
        "is_song_deck": looks_like_song_deck,
        "flags": {
            "SPARSE": sparse,
            "CROWDED": crowded,
            "TINY_TEXT": tiny,
            "TAIL": tail,
            "ORPHAN_START": orphan_start,
        },
        "slides": slide_stats,
        "issues": issues,
    }
