#!/usr/bin/env python3
"""
One-command dev quality gate.

Runs run_dev_tests.py, then evaluates PASS/FAIL using qa_report.json.
If FAIL, writes dev_out/qa/issue_bundle.json with details for fast iteration.

No network. No AI. Deterministic.
"""

from __future__ import annotations

import json
import subprocess
import sys
from pathlib import Path
from pptx import Presentation

PROJECT_DIR = Path(__file__).resolve().parent

DEFAULT_OUT_DIR = PROJECT_DIR / "dev_out"
DEFAULT_QA_JSON = DEFAULT_OUT_DIR / "qa" / "qa_report.json"
DEFAULT_BUNDLE = DEFAULT_OUT_DIR / "qa" / "issue_bundle.json"

# Tunable thresholds (feel free to tweak)
THRESHOLDS = {
    "songs": {
        "CROWDED": 0,
        "TINY_TEXT": 0,
        "SPARSE": 25,          # songs can legitimately have more sparse slides
        "TAIL": 999,           # not used for songs
        "ORPHAN_START": 999,
    },
    "verses": {
        "CROWDED": 0,
        "TINY_TEXT": 0,
        "SPARSE": 6,           # target to drive down over time
        "TAIL": 3,             # key metric: keep tails very low
        "ORPHAN_START": 3,     # key metric: avoid unnatural starts
    }
}

def _run(cmd: list[str]) -> int:
    proc = subprocess.run(cmd, cwd=str(PROJECT_DIR))
    return proc.returncode

def _extract_slide_texts(pptx_path: Path, slide_numbers: list[int]) -> dict[int, str]:
    """Return {slide_num: text} for selected slides."""
    prs = Presentation(str(pptx_path))
    out: dict[int, str] = {}
    wanted = set(slide_numbers)
    for idx, slide in enumerate(prs.slides, start=1):
        if idx not in wanted:
            continue
        parts = []
        for sh in slide.shapes:
            if getattr(sh, "has_text_frame", False):
                try:
                    t = (sh.text_frame.text or "").strip()
                except Exception:
                    t = ""
                if t:
                    parts.append(t)
        out[idx] = "\n---\n".join(parts).strip()
    return out

def _evaluate(report: dict) -> tuple[bool, dict]:
    """Return (pass, details)."""
    details = {}
    all_pass = True
    for section in ("songs", "verses"):
        r = report.get(section, {})
        flags = r.get("flags", {})
        sec_details = {"slide_count": r.get("slide_count"), "counts": {}, "fails": {}}
        for key, limit in THRESHOLDS[section].items():
            cnt = len(flags.get(key, []))
            sec_details["counts"][key] = cnt
            if cnt > limit:
                sec_details["fails"][key] = {"count": cnt, "limit": limit, "slides": flags.get(key, [])}
                all_pass = False
        details[section] = sec_details
    return all_pass, details

def main(argv: list[str]) -> int:
    # Forward args to run_dev_tests if user provided any
    run_cmd = [sys.executable, str(PROJECT_DIR / "run_dev_tests.py")] + argv[1:]
    rc = _run(run_cmd)
    if rc != 0:
        print(f"dev_check: run_dev_tests failed (exit {rc})")
        return rc

    if not DEFAULT_QA_JSON.exists():
        print(f"dev_check: missing {DEFAULT_QA_JSON}")
        return 2

    report = json.loads(DEFAULT_QA_JSON.read_text(encoding="utf-8"))
    ok, details = _evaluate(report)

    # Build bundle if failing (or always, for convenience)
    bundle = {
        "ok": ok,
        "thresholds": THRESHOLDS,
        "details": details,
        "paths": {
            "songs_pptx": report["songs"]["pptx"],
            "verses_pptx": report["verses"]["pptx"],
            "qa_json": str(DEFAULT_QA_JSON),
        },
        "slides": {"songs": {}, "verses": {}},
    }

    if not ok:
        # collect slide texts for the most actionable verse failures
        verses_flags = report["verses"]["flags"]
        focus = sorted(set(
            verses_flags.get("TAIL", []) +
            verses_flags.get("ORPHAN_START", []) +
            verses_flags.get("SPARSE", [])
        ))
        focus = focus[:40]
        bundle["slides"]["verses"] = _extract_slide_texts(Path(report["verses"]["pptx"]), focus)

    DEFAULT_BUNDLE.write_text(json.dumps(bundle, indent=2), encoding="utf-8")

    # Print summary
    print("\n=== DEV CHECK SUMMARY ===")
    for section in ("songs", "verses"):
        c = details[section]["counts"]
        print(f"{section.upper():6} slides={details[section]['slide_count']} | "
              f"sparse={c.get('SPARSE',0)} tail={c.get('TAIL',0)} orphan={c.get('ORPHAN_START',0)} "
              f"crowded={c.get('CROWDED',0)} tiny={c.get('TINY_TEXT',0)}")
    print("STATUS:", "PASS" if ok else "FAIL")
    if not ok:
        print(f"Wrote issue bundle: {DEFAULT_BUNDLE}")

    return 0 if ok else 3

if __name__ == "__main__":
    raise SystemExit(main(sys.argv))
