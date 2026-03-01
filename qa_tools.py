from __future__ import annotations

from pathlib import Path
from pptx import Presentation

EMU_PER_PT = 12700

_ORPHAN_START_WORDS = {
    "and","but","for","so","to","of","or","nor","yet","then","also","thus"
}

_REF_LINE_RE = __import__("re").compile(r"\b\d+\s*:\s*\d+(?:\s*[-–]\s*\d+)?\b")

def _strip_reference_line(lines: list[str]) -> list[str]:
    """Remove trailing scripture reference line like 'Daniel 9:24' if present."""
    if not lines:
        return lines
    last = lines[-1].strip()
    # Common formats: 'Daniel 9:24' or 'Matthew 28:19-20'
    if _REF_LINE_RE.search(last):
        return lines[:-1]
    return lines

def _ends_with_punct(lines: list[str]) -> bool:
    lines = [ln.strip() for ln in lines if ln.strip()]
    lines = _strip_reference_line(lines)
    if not lines:
        return False
    last = lines[-1].rstrip()
    return bool(last) and last[-1] in ".;:!?"

def _iter_text_shapes(slide):
    for sh in slide.shapes:
        if getattr(sh, "has_text_frame", False):
            try:
                txt = sh.text_frame.text
            except Exception:
                txt = ""
            if txt is not None:
                yield sh, txt

def analyze_pptx(pptx_path: Path) -> dict:
    """Lightweight QA heuristics for quickly spotting bad slides.

    Flags:
      - SPARSE: too little text (often indicates broken grouping)
      - TAIL: 1–2 line leftover that probably should merge
      - ORPHAN_START: slide starts with a weak connector (And/But/So/To/…) and is short
      - CROWDED: too much text (hard to read)
      - TINY_TEXT: font likely shrank (autosize or template override)
    """
    prs = Presentation(str(pptx_path))
    sparse = []
    tail = []
    orphan = []
    crowded = []
    tiny = []
    slide_stats = []

    prev_lines: list[str] = []

    for idx, slide in enumerate(prs.slides, start=1):
        texts = []
        min_font_pt = None
        for sh, txt in _iter_text_shapes(slide):
            t = (txt or "").strip()
            if not t:
                continue
            texts.append(t)
            # estimate min font size on first run of each paragraph
            try:
                for p in sh.text_frame.paragraphs:
                    if p.runs and p.runs[0].font.size:
                        pt = float(p.runs[0].font.size.pt)
                        min_font_pt = pt if min_font_pt is None else min(min_font_pt, pt)
                    elif p.font and p.font.size:
                        pt = float(p.font.size.pt)
                        min_font_pt = pt if min_font_pt is None else min(min_font_pt, pt)
            except Exception:
                pass

        full = "\n".join(texts).strip()

        # Ignore token-only template slides if present
        if "{{" in full and "}}" in full:
            slide_stats.append({"slide": idx, "ignored": True, "chars": 0, "lines": 0, "min_font_pt": min_font_pt})
            prev_full_text = full
            continue

        chars = len(full.replace("\n"," ").strip())
        lines = sum(1 for ln in full.splitlines() if ln.strip())
        slide_stats.append({"slide": idx, "chars": chars, "lines": lines, "min_font_pt": min_font_pt})

        # Heuristics (conservative defaults)
        if chars > 360 or lines > 8:
            crowded.append(idx)

        if 0 < chars < 45 or (0 < lines < 2):
            sparse.append(idx)

        # Tail: short, 1–2 line leftovers (avoid false positives on title/ref-only slides)
        if 0 < chars < 120 and 1 <= lines <= 2:
            tail.append(idx)

        # Orphan start: ONLY flag when it looks like a *bad split*, not a legitimate KJV verse opening.
        # Rule:
        #   - Slide starts with a connector word (And/But/So/For/…)
        #   - AND the previous slide did NOT end in punctuation (. ; : ! ?)
        # This avoids false positives on verses that genuinely begin with "And ...".
        if idx > 1 and 0 < chars < 220 and 1 <= lines <= 3:
            first_line = next((ln.strip() for ln in full.splitlines() if ln.strip()), "")
            first_word = first_line.split(" ", 1)[0].strip("“”'\"(),;:").lower() if first_line else ""
            if first_word in _ORPHAN_START_WORDS and not _ends_with_punct(prev_lines):
                orphan.append(idx)

        if min_font_pt is not None and min_font_pt < 26:
            tiny.append(idx)

        prev_lines = [ln for ln in full.splitlines() if ln.strip()]

    return {
        "pptx": str(pptx_path),
        "slide_count": len(prs.slides),
        "flags": {"SPARSE": sparse, "TAIL": tail, "ORPHAN_START": orphan, "CROWDED": crowded, "TINY_TEXT": tiny},
        "slides": slide_stats,
    }